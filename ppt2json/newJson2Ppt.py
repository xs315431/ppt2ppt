from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import json
import os


def hex_to_rgb(hex_str):
    if hex_str and len(hex_str) == 6:
        return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))
    return None


def add_textbox(slide, element):
    opts = element['options']
    textbox = slide.shapes.add_textbox(Inches(opts['x']), Inches(opts['y']), Inches(opts['w']), Inches(opts['h']))
    tf = textbox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = element['text']

    if opts.get('fontSize'):
        run.font.size = Pt(opts['fontSize'])
    if opts.get('bold') is not None:
        run.font.bold = opts['bold']
    if opts.get('italic') is not None:
        run.font.italic = opts['italic']
    if opts.get('color'):
        run.font.color.rgb = hex_to_rgb(opts['color'])
    if opts.get('align'):
        try:
            p.alignment = PP_ALIGN[opts['align'].upper()]
        except KeyError:
            pass


def add_image(slide, element, base_dir):
    opts = element['options']
    rel_path = element['path']
    abs_path = os.path.join(base_dir, rel_path)
    if not os.path.exists(abs_path):
        print(f"⚠️ 图片路径不存在：{abs_path}")
        return
    slide.shapes.add_picture(abs_path, Inches(opts['x']), Inches(opts['y']), width=Inches(opts['w']), height=Inches(opts['h']))


def add_shape(slide, element):
    opts = element['options']
    shape_type = getattr(MSO_AUTO_SHAPE_TYPE, element.get('shape_type', 'RECTANGLE'), MSO_AUTO_SHAPE_TYPE.RECTANGLE)
    shape = slide.shapes.add_shape(shape_type, Inches(opts['x']), Inches(opts['y']), Inches(opts['w']), Inches(opts['h']))
    shape.text = element.get("text", "")


def add_table(slide, element):
    opts = element['options']
    rows = element['rows']
    cols = element['cols']
    data = element['data']

    table_shape = slide.shapes.add_table(rows, cols, Inches(opts['x']), Inches(opts['y']), Inches(opts['w']), Inches(opts['h']))
    table = table_shape.table

    for r in range(rows):
        for c in range(cols):
            text = data[r][c] if r < len(data) and c < len(data[r]) else ""
            table.cell(r, c).text = text


def json_to_pptx(json_path, output_path):
    base_dir = os.path.dirname(json_path)
    with open(json_path, 'r', encoding='utf-8') as f:
        json_data = json.load(f)

    prs = Presentation()
    blank_layout = prs.slide_layouts[6]  # 空白布局

    for slide_data in json_data['slides']:
        slide = prs.slides.add_slide(blank_layout)
        for element in slide_data['elements']:
            try:
                if element['type'] == 'text':
                    add_textbox(slide, element)
                elif element['type'] == 'image':
                    add_image(slide, element, base_dir)
                elif element['type'] == 'shape':
                    add_shape(slide, element)
                elif element['type'] == 'table':
                    add_table(slide, element)
                elif element['type'] == 'chart':
                    print("⚠️ 暂不支持 chart 元素还原，跳过")
                else:
                    print(f"⚠️ 未知元素类型 {element['type']}，跳过")
            except Exception as e:
                print(f"❌ 处理元素时出错：{e}")

    prs.save(output_path)
    print(f"✅ PPTX 还原完成，保存为：{output_path}")


# 示例用法
if __name__ == "__main__":
    json_input_path = "./pptjson/output.json"
    pptx_output_path = "./ppt/restored.pptx"
    json_to_pptx(json_input_path, pptx_output_path)
