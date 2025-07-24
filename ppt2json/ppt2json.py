from pptx import Presentation
import json
import os

# 需要处理所有类型的文件，将它存储为json文件配置
def extract_ppt_to_json(ppt_path: str) -> dict:
    prs = Presentation(ppt_path)
    slides_json = []

    def pt_to_inches(pt):
        return round(pt / 72, 2)

    for slide_idx, slide in enumerate(prs.slides):
        slide_data = {
            "title": f"Slide {slide_idx + 1}",
            "backgroundColor": "",
            "backgroundImage": "",
            "elements": []
        }

        for shape in slide.shapes:
            # 处理图片
            if shape.shape_type == 13:  # PICTURE
                image_path = f"image_slide{slide_idx + 1}_{len(slide_data['elements'])}.png"
                with open(image_path, 'wb') as f:
                    f.write(shape.image.blob)

                slide_data["elements"].append({
                    "type": "image",
                    "path": image_path,
                    "options": {
                        "path": image_path,
                        "x": pt_to_inches(shape.left.pt),
                        "y": pt_to_inches(shape.top.pt),
                        "w": pt_to_inches(shape.width.pt),
                        "h": pt_to_inches(shape.height.pt)
                    }
                })

            # 处理文本
            elif shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    font_size = None
                    font_color = None
                    is_bold = False

                    for run in para.runs:
                        if not font_size and run.font.size:
                            font_size = run.font.size.pt
                        if not font_color and run.font.color and run.font.color.rgb:
                            font_color = str(run.font.color.rgb)
                        if run.font.bold:
                            is_bold = True

                    if not font_size:
                        font_size = 16  # 默认值
                    if not font_color:
                        font_color = "000000"

                    slide_data["elements"].append({
                        "type": "text",
                        "text": para.text,
                        "options": {
                            "x": pt_to_inches(shape.left.pt),
                            "y": pt_to_inches(shape.top.pt),
                            "w": pt_to_inches(shape.width.pt),
                            "h": pt_to_inches(shape.height.pt),
                            "fontSize": font_size,
                            "align": para.alignment.name.lower() if para.alignment else "left",
                            "color": font_color,
                            "bold": is_bold
                        }
                    })

        slides_json.append(slide_data)

    return {"slides": slides_json}

# 用法示例
ppt_file = "./PPTtemplate/temp1.pptx"
json_output = extract_ppt_to_json(ppt_file)

with open("./pptjson/output.json", "w", encoding="utf-8") as f:
    json.dump(json_output, f, ensure_ascii=False, indent=4)

print("✅ PPT 转 JSON 成功")
