from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
import json
import os
from pathlib import Path


def rgb_to_hex(rgb_color: RGBColor):
    if rgb_color:
        return '{:02X}{:02X}{:02X}'.format(rgb_color[0], rgb_color[1], rgb_color[2])
    return None


def extract_text_frame(text_frame):
    para = text_frame.paragraphs[0]
    run = para.runs[0] if para.runs else None
    font = run.font if run else None

    return {
        "text": text_frame.text,
        "fontSize": font.size.pt if font and font.size else None,
        "bold": font.bold if font else None,
        "italic": font.italic if font else None,
        "color": rgb_to_hex(font.color.rgb) if font and font.color and font.color.rgb else None,
        "align": para.alignment.name.lower() if para.alignment else None,
    }


def extract_shape(shape, slide_idx, shape_idx, image_output_dir):
    base = {
        "x": round(shape.left.inches, 2),
        "y": round(shape.top.inches, 2),
        "w": round(shape.width.inches, 2),
        "h": round(shape.height.inches, 2),
    }

    if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX or (shape.has_text_frame and shape.text.strip()):
        text_info = extract_text_frame(shape.text_frame)
        return {
            "type": "text",
            "text": text_info["text"],
            "options": {**base, **{k: v for k, v in text_info.items() if k != "text"}}
        }

    elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        # 提取图片并保存
        ext = shape.image.ext  # 如 '.png'
        img_bytes = shape.image.blob
        img_name = f"image_slide{slide_idx}_{shape_idx}.{ext}"
        img_path = os.path.join(image_output_dir, img_name)
        with open(img_path, 'wb') as f:
            f.write(img_bytes)

        return {
            "type": "image",
            "path": os.path.join("images", img_name),
            "options": base
        }

    elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
        return {
            "type": "shape",
            "shape_type": shape.auto_shape_type.name,
            "text": shape.text,
            "options": base
        }

    elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        table = shape.table
        rows, cols = len(table.rows), len(table.columns)
        cells = [[table.cell(r, c).text for c in range(cols)] for r in range(rows)]
        return {
            "type": "table",
            "rows": rows,
            "cols": cols,
            "data": cells,
            "options": base
        }

    elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
        chart = shape.chart
        chart_data = {
            "title": chart.chart_title.text_frame.text if chart.has_title else "",
            "has_legend": chart.has_legend,
            "chart_type": chart.chart_type.name
        }
        return {
            "type": "chart",
            "options": base,
            "chart": chart_data
        }

    else:
        return {
            "type": "unknown",
            "shape_type": shape.shape_type,
            "options": base
        }


def extract_pptx_to_json(pptx_path: str, output_json_path: str):
    prs = Presentation(pptx_path)
    slides_json = []

    # 准备图片目录
    output_dir = os.path.dirname(output_json_path)
    image_output_dir = os.path.join(output_dir, "images")
    os.makedirs(image_output_dir, exist_ok=True)

    for slide_idx, slide in enumerate(prs.slides):
        slide_elements = []
        for shape_idx, shape in enumerate(slide.shapes):
            try:
                element = extract_shape(shape, slide_idx, shape_idx, image_output_dir)
                if element:
                    slide_elements.append(element)
            except Exception as e:
                print(f"❌ Error extracting shape on slide {slide_idx}, shape {shape_idx}: {e}")
        slides_json.append({"elements": slide_elements})

    result = {"slides": slides_json}
    with open(output_json_path, 'w', encoding='utf-8') as f:
        json.dump(result, f, indent=2, ensure_ascii=False)

    print(f"✅ PPTX 提取完成，已保存 JSON 到 {output_json_path}，图片保存到 images/")


# 示例用法
if __name__ == "__main__":
    pptx_file = "./PPTTemplate/电商年中总结.pptx"
    output_json = "./pptjson/output.json"
    extract_pptx_to_json(pptx_file, output_json)
